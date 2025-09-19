"""
PPT 处理工具函数集合
包含 Markdown 解析、文本格式保持、结构化读取等功能
"""

import os
import re
import json
from typing import List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
    pptx_available = True
except ImportError:
    pptx_available = False


def read_text(path: str) -> str:
    """读取文本文件"""
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def parse_md_for_ppt_structure(md_text: str) -> dict:
    """
    解析通用 Markdown（以 /md_input_file/pptx_test_case.md 为参照）为结构化信息：
    - 标题页：title, org, date
    - 目录：toc(list[str])，来自"## 目录"节的有序列表
    - 正文：body_sections(list[dict])，过滤掉"目录/下一步行动/成功指标"等特殊节
      每个节包含：heading(str)，blocks(list[dict])，block.type in {paragraph,list,table}
    - 结尾：ending(dict) 包含 next_actions(list[str]) 与 success_metrics(list[str])
    """
    lines = md_text.splitlines()

    def is_heading(line: str, level: int) -> bool:
        prefix = "#" * level + " "
        return line.startswith(prefix)

    def strip_heading_marker(line: str) -> str:
        return re.sub(r"^#+\s+", "", line).strip()

    # 1) 标题页（前三个标题行）
    title = ""
    org = ""
    date = ""
    for i, line in enumerate(lines[:6]):
        if not title and is_heading(line, 1):
            title = strip_heading_marker(line)
            continue
        if not org and is_heading(line, 2):
            org = strip_heading_marker(line)
            continue
        if not date and is_heading(line, 3):
            date = strip_heading_marker(line)
            continue

    # 2) 遍历收集各节
    toc: List[str] = []
    body_sections: List[dict] = []
    ending = {"next_actions": [], "success_metrics": []}

    i = 0
    n = len(lines)

    def collect_list(start_idx: int) -> Tuple[List[str], int]:
        items: List[str] = []
        j = start_idx
        list_pattern = re.compile(r"^(?:[-*]\s+|\d+\.|\d+\)\s+)")
        while j < n:
            ln = lines[j].rstrip()
            if not ln:
                break
            if list_pattern.match(ln):
                # 去掉前缀
                item = list_pattern.sub("", ln).strip()
                items.append(item)
                j += 1
            else:
                break
        return items, j

    def collect_table(start_idx: int) -> Tuple[List[str], int]:
        j = start_idx
        tbl_lines: List[str] = []
        while j < n:
            ln = lines[j].rstrip()
            if ln.startswith("|") and "|" in ln:
                tbl_lines.append(ln)
                j += 1
            else:
                break
        return tbl_lines, j

    def collect_paragraph(start_idx: int) -> Tuple[str, int]:
        j = start_idx
        buff: List[str] = []
        while j < n:
            ln = lines[j].rstrip()
            if not ln:
                if buff:
                    j += 1
                    break
                else:
                    j += 1
                    continue
            if ln.startswith("#") or ln.startswith("---") or ln.startswith("|") or re.match(r"^(?:[-*]\s+|\d+\.|\d+\)\s+)", ln):
                break
            buff.append(ln)
            j += 1
        return "\n".join(buff).strip(), j

    current_section: Optional[dict] = None

    while i < n:
        line = lines[i].rstrip()

        # 跳过分割线
        if line.strip() == "---":
            i += 1
            continue

        # 捕获目录节：支持“目录”以及“目录（...）”形式
        if is_heading(line, 2) and strip_heading_marker(line).startswith("目录"):
            i += 1
            # 收集有序列表
            items, i = collect_list(i)
            toc.extend(items)
            continue

        # 结尾：下一步行动 / 成功指标
        if is_heading(line, 2) and strip_heading_marker(line) in ("下一步行动", "成功指标"):
            heading = strip_heading_marker(line)
            i += 1
            items, i = collect_list(i)
            if heading == "下一步行动":
                ending["next_actions"].extend(items)
            else:
                ending["success_metrics"].extend(items)
            continue

        # 正文节（二级标题，排除上面几类）
        if is_heading(line, 2):
            heading = strip_heading_marker(line)
            current_section = {"heading": heading, "blocks": []}
            body_sections.append(current_section)
            i += 1
            continue

        # 在某个正文节内，解析块
        if current_section is not None:
            if not line:
                i += 1
                continue

            # 分隔线（直接跳过并前进）
            if line.startswith("---"):
                i += 1
                continue

            # 小标题（### 等）：作为一个独立块消费掉，避免死循环
            if line.startswith("#"):
                current_section["blocks"].append({
                    "type": "subheading",
                    "text": strip_heading_marker(line)
                })
                i += 1
                continue

            # 列表
            if re.match(r"^(?:[-*]\s+|\d+\.|\d+\)\s+)", line):
                items, i = collect_list(i)
                if items:
                    current_section["blocks"].append({"type": "list", "items": items})
                continue

            # 表格
            if line.startswith("|") and "|" in line:
                tbl_lines, i = collect_table(i)
                if tbl_lines:
                    current_section["blocks"].append({"type": "table", "lines": tbl_lines})
                continue

            # 段落
            para, i2 = collect_paragraph(i)
            if para:
                current_section["blocks"].append({"type": "paragraph", "text": para})
                i = i2
            else:
                # 未能收集到段落，至少前进一步避免卡死
                i += 1
            continue

        i += 1

    return {
        "title_page": {"title": title, "org": org, "date": date},
        "toc": toc,
        "body_sections": body_sections,
        "ending": ending,
    }


def update_text_preserve_format(shape, new_text):
    """更新文本但保持原有格式，支持多行文本"""
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        return
    
    # 分割新文本为多行
    lines = new_text.split('\n')
    
    # 保存原有段落的格式
    original_paragraphs = []
    for para in text_frame.paragraphs:
        if para.runs:
            original_font = para.runs[0].font
            original_paragraphs.append({
                'font_name': original_font.name,
                'font_size': original_font.size,
                'font_bold': original_font.bold,
                'font_italic': original_font.italic,
                'font_color': original_font.color,
                'alignment': para.alignment
            })
        else:
            original_paragraphs.append(None)
    
    # 清空所有段落
    for para in text_frame.paragraphs:
        para.clear()
    
    # 确保有足够的段落
    while len(text_frame.paragraphs) < len(lines):
        text_frame.add_paragraph()
    
    # 为每行文本设置格式
    for i, line in enumerate(lines):
        if i < len(text_frame.paragraphs):
            para = text_frame.paragraphs[i]
            new_run = para.add_run()
            new_run.text = line.strip()
            
            # 使用对应段落的原有格式，如果没有则使用第一个段落的格式
            format_index = min(i, len(original_paragraphs) - 1)
            if format_index >= 0 and original_paragraphs[format_index]:
                orig_format = original_paragraphs[format_index]
                
                # 恢复字体格式
                new_run.font.name = orig_format['font_name']
                if orig_format['font_size']:
                    new_run.font.size = orig_format['font_size']
                new_run.font.bold = orig_format['font_bold']
                new_run.font.italic = orig_format['font_italic']
                
                # 恢复段落对齐
                if orig_format['alignment']:
                    para.alignment = orig_format['alignment']
                
                # 恢复颜色
                try:
                    original_color = orig_format['font_color']
                    if hasattr(original_color, 'rgb') and original_color.rgb:
                        new_run.font.color.rgb = original_color.rgb
                    elif hasattr(original_color, 'theme_color') and original_color.theme_color:
                        new_run.font.color.theme_color = original_color.theme_color
                except Exception:
                    pass


def get_shape_info(shape) -> dict:
    """获取形状的详细信息"""
    info = {
        "name": getattr(shape, "name", ""),
        "has_text": hasattr(shape, "text_frame"),
        "text": getattr(shape, "text", None) if hasattr(shape, "text") else None,
        "is_placeholder": getattr(shape, "is_placeholder", False),
        "placeholder_type": None,
        "placeholder_idx": None,
        "shape_type": getattr(getattr(shape, "shape_type", None), "name", None),
        "left_in": float(shape.left.inches) if hasattr(shape, "left") else None,
        "top_in": float(shape.top.inches) if hasattr(shape, "top") else None,
        "width_in": float(shape.width.inches) if hasattr(shape, "width") else None,
        "height_in": float(shape.height.inches) if hasattr(shape, "height") else None,
    }
    try:
        if info["is_placeholder"] and pptx_available:
            pf = shape.placeholder_format
            info["placeholder_type"] = getattr(getattr(pf, "type", None), "name", None)
            info["placeholder_idx"] = getattr(pf, "idx", None)
    except Exception:
        pass
    return info


def extract_sections_from_md(md_text: str) -> dict:
    """旧版简单提取函数（兼容性保留）"""
    data = {"title": "", "org": "", "date": "", "toc": [], "para": "", "image": None, "table": None}
    lines = [l.strip() for l in md_text.splitlines()]
    i = 0
    while i < len(lines):
        line = lines[i]
        if line.startswith("- **总标题**:"):
            data["title"] = line.split(":", 1)[1].strip()
        elif line.startswith("- **制作单位**:"):
            data["org"] = line.split(":", 1)[1].strip()
        elif line.startswith("- **时间**:"):
            data["date"] = line.split(":", 1)[1].strip()
        elif line.startswith("- **目录**:"):
            i += 1
            while i < len(lines) and lines[i].startswith("- "):
                item = lines[i][2:].strip()
                if item:
                    data["toc"].append(item)
                i += 1
            i -= 1
        elif line.startswith("- **正文页 1（段落 + 图片）**:"):
            i += 1
            while i < len(lines) and lines[i].startswith("- "):
                sub = lines[i]
                if sub.startswith("- 段落:"):
                    data["para"] = sub.split(":", 1)[1].strip()
                if sub.startswith("- 图片:"):
                    data["image"] = sub.split(":", 1)[1].strip() or None
                i += 1
            i -= 1
        elif line.startswith("- **正文页 2（表格）**:"):
            data["table"] = {
                "rows": 4,
                "cols": 4,
                "data": [
                    ["项目", "人工成本(周)", "自动化成本(周)", "差值(周)"],
                    ["周报生成", "2.0", "0.2", "1.8"],
                    ["投标材料", "3.5", "1.0", "2.5"],
                    ["数据汇总", "1.0", "0.3", "0.7"],
                ],
            }
        i += 1
    return data


def resolve_path(path: str) -> str:
    """解析路径"""
    if os.path.isabs(path):
        return path
    base_path = os.environ.get('OFFICE_EDIT_PATH')
    if not base_path:
        base_path = os.path.join(os.path.expanduser('~'), '桌面')
    return os.path.join(base_path, path)


def resolve_template_path(template_path: str) -> str:
    """解析模板路径，支持容错匹配"""
    if os.path.exists(template_path):
        return template_path
    dir_name = os.path.dirname(template_path)
    base_try = os.path.basename(template_path).replace(" ", "")
    candidates = []
    try:
        for name in os.listdir(dir_name or "."):
            if name.lower().endswith(".pptx"):
                norm = name.replace(" ", "")
                if norm.lower() == base_try.lower():
                    candidates.append(os.path.join(dir_name, name))
        if not candidates:
            # 再退一步：拿目录下任意一个 .pptx 作为模板
            for name in os.listdir(dir_name or "."):
                if name.lower().endswith(".pptx"):
                    candidates.append(os.path.join(dir_name, name))
        if candidates:
            return candidates[0]
    except Exception:
        pass
    return template_path


def split_title_by_length(title_text: str, first_ratio: float = 0.4) -> str:
    """根据比例自动将标题拆分为两行: 第一行40%，第二行60%（字母和空格算0.5个字符宽度）"""
    if not title_text:
        return ""
    
    def get_char_width(char):
        """获取字符的视觉宽度：字母和空格为0.5，中文字符为1.0"""
        if char.isascii() and (char.isalpha() or char == ' '):
            return 0.5
        else:
            return 1.0
    
    def get_text_width(text):
        """计算文本的总宽度"""
        return sum(get_char_width(char) for char in text)
    
    # 计算总宽度和目标分割宽度
    total_width = get_text_width(title_text)
    target_width = total_width * first_ratio
    
    # 寻找最佳分割点
    best_split = 0
    best_width_diff = float('inf')
    found_good_split = False
    
    # 遍历所有可能的分割点
    current_width = 0
    for i in range(len(title_text)):
        current_width += get_char_width(title_text[i])
        
        # 计算当前位置的宽度差异
        width_diff = abs(current_width - target_width)
        
        # 如果在空格处，优先选择
        if title_text[i] == ' ':
            if width_diff < target_width * 0.5:  # 在目标宽度50%范围内
                best_split = i
                found_good_split = True
                break
        
        # 如果在其他分割符处，且更接近目标
        elif title_text[i] in ['，', '。', '：', '；', '、', '－', '-']:
            if width_diff < best_width_diff:
                best_split = i + 1
                best_width_diff = width_diff
                found_good_split = True
        
        # 如果没找到分割符，记录最接近目标宽度的位置
        elif not found_good_split and width_diff < best_width_diff:
            best_split = i
            best_width_diff = width_diff
    
    # 处理空格分割的情况
    if best_split < len(title_text) and title_text[best_split] == ' ':
        first_line = title_text[:best_split].strip()
        second_line = title_text[best_split+1:].strip()
    else:
        first_line = title_text[:best_split].strip()
        second_line = title_text[best_split:].strip()
    
    # 调试输出
    first_width = get_text_width(first_line)
    second_width = get_text_width(second_line)
    print(f"原标题: '{title_text}' (总宽度: {total_width:.1f})")
    print(f"目标分割宽度: {target_width:.1f} (40%)")
    print(f"实际分割位置: {best_split}")
    print(f"第一行: '{first_line}' (宽度: {first_width:.1f})")
    print(f"第二行: '{second_line}' (宽度: {second_width:.1f})")
    print(f"比例: {first_width/total_width*100:.1f}% : {second_width/total_width*100:.1f}%")
    
    return f"{first_line}\n{second_line}"


def smart_update_toc_items(slide, toc_items: list):
    """智能更新目录项，支持2-5个目录的自动匹配，保持原有模板布局美观性"""
    if not pptx_available:
        return
    
    # 模板原有的3个目录项位置和形状索引
    original_toc_mapping = [
        {'content_idx': 3, 'number_idx': 4},  # 目录标题1, 01
        {'content_idx': 5, 'number_idx': 7},  # 目录标题2, 02  
        {'content_idx': 6, 'number_idx': 8}   # 目录标题3, 03
    ]
    
    # 限制目录项数量在2-6之间（模板适配双列最多3行/列）
    toc_count = max(2, min(6, len(toc_items)))
    actual_items = toc_items[:toc_count]
    
    print(f"智能匹配目录：{len(toc_items)} 项 → {toc_count} 项")
    
    try:
        # 策略1：如果是2-3个目录项，直接更新现有位置，清空多余位置
        if toc_count <= 3:
            for i in range(3):
                if i < len(original_toc_mapping):
                    mapping = original_toc_mapping[i]
                    
                    # 更新内容
                    if mapping['content_idx'] < len(slide.shapes):
                        if i < toc_count:
                            update_text_preserve_format(slide.shapes[mapping['content_idx']], actual_items[i])
                        else:
                            # 清空多余的位置
                            update_text_preserve_format(slide.shapes[mapping['content_idx']], "")
                    
                    # 更新编号
                    if mapping['number_idx'] < len(slide.shapes):
                        if i < toc_count:
                            update_text_preserve_format(slide.shapes[mapping['number_idx']], f"{i+1:02d}")
                        else:
                            # 清空多余的编号
                            update_text_preserve_format(slide.shapes[mapping['number_idx']], "")
        
        # 策略2：如果是4-6个目录项，使用双列交替布局
        else:
            # 先清空所有原有目录项
            for i in range(3):
                mapping = original_toc_mapping[i]
                if mapping['content_idx'] < len(slide.shapes):
                    update_text_preserve_format(slide.shapes[mapping['content_idx']], "")
                if mapping['number_idx'] < len(slide.shapes):
                    update_text_preserve_format(slide.shapes[mapping['number_idx']], "")
            
            # 获取模板形状的格式参数
            content_template = slide.shapes[original_toc_mapping[0]['content_idx']] if original_toc_mapping[0]['content_idx'] < len(slide.shapes) else None
            number_template = slide.shapes[original_toc_mapping[0]['number_idx']] if original_toc_mapping[0]['number_idx'] < len(slide.shapes) else None
            
            # 双列交替布局参数
            left_content_left = 3.7
            left_number_left = 2.5
            right_content_left = left_content_left + 4.5  # 右移4.5英寸
            right_number_left = left_number_left + 4.5
            
            base_top = 2.5
            number_offset = -0.1  # 编号稍微向上偏移
            vertical_spacing = 1.05  # 稍微压缩竖向间距，容纳3行
            
            print(f"使用双列交替布局：左列奇数项，右列偶数项")
            
            # 为所有目录项创建新文本框（双列交替）
            for i in range(toc_count):
                item_num = i + 1
                is_odd = (item_num % 2 == 1)  # 奇数项放左列，偶数项放右列
                
                if is_odd:
                    # 奇数项：左列 (01, 03, 05)
                    content_left = left_content_left
                    number_left = left_number_left
                    row_in_column = (item_num - 1) // 2  # 在左列中的行号：0, 1, 2
                else:
                    # 偶数项：右列 (02, 04)
                    content_left = right_content_left
                    number_left = right_number_left
                    row_in_column = (item_num - 2) // 2  # 在右列中的行号：0, 1
                
                content_top = base_top + row_in_column * vertical_spacing
                number_top = content_top + number_offset
                
                column_name = "左列" if is_odd else "右列"
                print(f"目录项{item_num:02d}({column_name}): 内容({content_left:.1f}, {content_top:.1f}) 编号({number_left:.1f}, {number_top:.1f})")
                
                # 创建内容文本框（每列最多3项 -> 行号最多为2）
                content_textbox = slide.shapes.add_textbox(
                    Inches(content_left),
                    Inches(content_top),
                    Inches(4.0),  # 稍微缩小宽度以适应双列
                    Inches(0.6)
                )
                
                # 设置内容文本并复制格式
                content_frame = content_textbox.text_frame
                content_frame.clear()
                content_para = content_frame.paragraphs[0]
                content_run = content_para.add_run()
                content_run.text = actual_items[i]
                
                # 复制格式
                if content_template and content_template.text_frame.paragraphs[0].runs:
                    template_font = content_template.text_frame.paragraphs[0].runs[0].font
                    content_run.font.name = template_font.name
                    if template_font.size:
                        content_run.font.size = template_font.size
                    content_run.font.bold = template_font.bold
                    content_run.font.italic = template_font.italic
                    
                    try:
                        if hasattr(template_font.color, 'rgb') and template_font.color.rgb:
                            content_run.font.color.rgb = template_font.color.rgb
                        elif hasattr(template_font.color, 'theme_color') and template_font.color.theme_color:
                            content_run.font.color.theme_color = template_font.color.theme_color
                    except Exception:
                        pass
                
                # 创建编号文本框
                number_textbox = slide.shapes.add_textbox(
                    Inches(number_left),
                    Inches(number_top),
                    Inches(0.9),
                    Inches(0.8)
                )
                
                # 设置编号文本并复制格式
                number_frame = number_textbox.text_frame
                number_frame.clear()
                number_para = number_frame.paragraphs[0]
                number_run = number_para.add_run()
                number_run.text = f"{item_num:02d}"
                
                # 复制格式
                if number_template and number_template.text_frame.paragraphs[0].runs:
                    template_font = number_template.text_frame.paragraphs[0].runs[0].font
                    number_run.font.name = template_font.name
                    if template_font.size:
                        number_run.font.size = template_font.size
                    number_run.font.bold = template_font.bold
                    number_run.font.italic = template_font.italic
                    
                    try:
                        if hasattr(template_font.color, 'rgb') and template_font.color.rgb:
                            number_run.font.color.rgb = template_font.color.rgb
                        elif hasattr(template_font.color, 'theme_color') and template_font.color.theme_color:
                            number_run.font.color.theme_color = template_font.color.theme_color
                    except Exception:
                        pass
    
    except Exception as e:
        print(f"智能更新目录时出错: {e}")


def extract_subsection_content_from_md(md_text: str) -> list:
    """从Markdown文本中提取子章节内容，每个子标题(###)单独成页"""
    if not md_text:
        return []
    
    lines = md_text.split('\n')
    subsections = []
    current_chapter_number = 0
    current_chapter_title = ""
    current_subsection = None
    
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        
        # 跳过前置内容（标题、目录等）
        if line.startswith('# ') or line.startswith('## 目录') or line == '---':
            i += 1
            continue
            
        # 检测章节标题 (## 1. 背景与目标)
        if line.startswith('## ') and any(char.isdigit() for char in line[:10]):
            # 提取章节编号
            import re
            match = re.match(r'## (\d+)\.\s*(.+)', line)
            if match:
                current_chapter_number = int(match.group(1))
                current_chapter_title = match.group(2).strip()
            # 每遇到新章节，重置当前子章节
            if current_subsection:
                subsections.append(current_subsection)
                current_subsection = None
            i += 1
            continue
        
        # 检测子标题 (### 1.1 办公痛点驱动自动化需求)
        elif line.startswith('### ') and current_chapter_number > 0:
            # 保存之前的子章节
            if current_subsection:
                subsections.append(current_subsection)
            
            # 开始新的子章节
            current_subsection = {
                'chapter_number': current_chapter_number,
                'chapter_title': current_chapter_title,  # 添加章节标题
                'title': line[4:].strip(),  # 去掉 "### "
                'content_blocks': []
            }
        
        # 收集子章节内容
        elif current_subsection and line:
            # 跳过下一个章节标题
            if line.startswith('## ') and any(char.isdigit() for char in line[:10]):
                # 遇到新章节，不要递增i，让外层循环处理
                continue
            
            # 表格开始
            if '|' in line and i+1 < len(lines) and '-' in lines[i+1]:
                table_lines = [line]
                i += 1
                # 收集表格头分隔行
                if i < len(lines):
                    table_lines.append(lines[i])
                    i += 1
                # 收集表格数据行
                while i < len(lines) and '|' in lines[i]:
                    table_lines.append(lines[i])
                    i += 1
                i -= 1  # 回退一行，因为最后会 i += 1
                
                # 为表格生成标题
                table_caption = f"表格 {len([b for b in current_subsection['content_blocks'] if b['type'] == 'table']) + 1}"
                current_subsection['content_blocks'].append({
                    'type': 'table',
                    'lines': table_lines,
                    'caption': table_caption
                })
            
            # 图片
            elif re.match(r'^!\[[^\]]*\]\([^\)]+\)', line):
                m = re.match(r'^!\[([^\]]*)\]\(([^\)]+)\)', line)
                if m:
                    alt_text = m.group(1).strip()  # 图片的alt文本作为标题
                    src = m.group(2).strip()
                    current_subsection['content_blocks'].append({
                        'type': 'image',
                        'src': src,
                        'caption': alt_text if alt_text else f"图片 {len([b for b in current_subsection['content_blocks'] if b['type'] == 'image']) + 1}"
                    })
            
            # 列表项
            elif line.startswith('- '):
                # 收集连续的列表项
                list_items = [line[2:].strip()]
                j = i + 1
                while j < len(lines) and lines[j].strip().startswith('- '):
                    list_items.append(lines[j].strip()[2:].strip())
                    j += 1
                i = j - 1  # 调整索引
                
                current_subsection['content_blocks'].append({
                    'type': 'list',
                    'items': list_items
                })
            
            # 普通段落
            elif not line.startswith('#') and line != '---':
                current_subsection['content_blocks'].append({
                    'type': 'paragraph',
                    'text': line
                })
        
        # 章节内的直系内容（没有任何 ### 子标题时）
        elif current_chapter_number > 0 and line:
            # 如遇到章节内直接出现的内容（图片、表格、段落、列表），为其创建一个隐式子章节
            if current_subsection is None:
                implicit_title = f"{current_chapter_number}.0 {current_chapter_title or '概览'}"
                current_subsection = {
                    'chapter_number': current_chapter_number,
                    'title': implicit_title,
                    'content_blocks': []
                }
            # 表格
            if '|' in line and i+1 < len(lines) and '-' in lines[i+1]:
                table_lines = [line]
                i += 1
                if i < len(lines):
                    table_lines.append(lines[i])
                    i += 1
                while i < len(lines) and '|' in lines[i]:
                    table_lines.append(lines[i])
                    i += 1
                i -= 1
                current_subsection['content_blocks'].append({'type': 'table', 'lines': table_lines})
            # 图片
            elif re.match(r'^!\[[^\]]*\]\([^\)]+\)', line):
                m = re.match(r'^!\[([^\]]*)\]\(([^\)]+)\)', line)
                if m:
                    alt_text = m.group(1).strip()  # 图片的alt文本作为标题
                    src = m.group(2).strip()
                    current_subsection['content_blocks'].append({
                        'type': 'image',
                        'src': src,
                        'caption': alt_text if alt_text else f"图片 {len([b for b in current_subsection['content_blocks'] if b['type'] == 'image']) + 1}"
                    })
            # 列表
            elif line.startswith('- '):
                list_items = [line[2:].strip()]
                j = i + 1
                while j < len(lines) and lines[j].strip().startswith('- '):
                    list_items.append(lines[j].strip()[2:].strip())
                    j += 1
                i = j - 1
                current_subsection['content_blocks'].append({'type': 'list', 'items': list_items})
            # 段落
            elif not line.startswith('#') and line != '---':
                current_subsection['content_blocks'].append({'type': 'paragraph', 'text': line})
        
        i += 1
    
    # 保存最后一个子章节
    if current_subsection:
        subsections.append(current_subsection)
    
    return subsections


def extract_chapter_content_from_md(md_text: str) -> list:
    """从Markdown文本中提取章节内容，每个章节包含标题和正文"""
    if not md_text:
        return []
    
    lines = md_text.split('\n')
    chapters = []
    current_chapter = None
    
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        
        # 跳过前置内容（标题、目录等）
        if line.startswith('# ') or line.startswith('## 目录') or line == '---':
            i += 1
            continue
            
        # 检测章节标题 (## 1. 背景与目标)
        if line.startswith('## ') and any(char.isdigit() for char in line[:10]):
            # 保存之前的章节
            if current_chapter:
                chapters.append(current_chapter)
            
            # 开始新章节
            current_chapter = {
                'title': line[3:].strip(),  # 去掉 "## "
                'content_blocks': []
            }
        
        # 收集章节内容
        elif current_chapter and line:
            # 子标题
            if line.startswith('### '):
                current_chapter['content_blocks'].append({
                    'type': 'subtitle',
                    'text': line[4:].strip()
                })
            
            # 表格开始
            elif '|' in line and '-' in lines[i+1] if i+1 < len(lines) else False:
                table_lines = [line]
                i += 1
                # 收集表格头分隔行
                if i < len(lines):
                    table_lines.append(lines[i])
                    i += 1
                # 收集表格数据行
                while i < len(lines) and '|' in lines[i]:
                    table_lines.append(lines[i])
                    i += 1
                i -= 1  # 回退一行，因为最后会 i += 1
                
                current_chapter['content_blocks'].append({
                    'type': 'table',
                    'lines': table_lines
                })
            
            # 列表项
            elif line.startswith('- '):
                # 收集连续的列表项
                list_items = [line[2:].strip()]
                j = i + 1
                while j < len(lines) and lines[j].strip().startswith('- '):
                    list_items.append(lines[j].strip()[2:].strip())
                    j += 1
                i = j - 1  # 调整索引
                
                current_chapter['content_blocks'].append({
                    'type': 'list',
                    'items': list_items
                })
            
            # 普通段落
            elif not line.startswith('#') and line != '---':
                current_chapter['content_blocks'].append({
                    'type': 'paragraph',
                    'text': line
                })
        
        i += 1
    
    # 保存最后一个章节
    if current_chapter:
        chapters.append(current_chapter)
    
    return chapters


def get_slide_master_layouts(presentation_path: str) -> dict:
    """获取演示文稿的母版布局信息"""
    if not pptx_available:
        return {"error": "python-pptx not available"}
    
    try:
        prs = Presentation(presentation_path)
        layouts_info = []
        
        for i, layout in enumerate(prs.slide_layouts):
            layout_info = {
                "index": i,
                "name": layout.name,
                "placeholders": []
            }
            
            # 获取布局中的占位符信息
            for j, placeholder in enumerate(layout.placeholders):
                try:
                    ph_info = {
                        "idx": placeholder.placeholder_format.idx,
                        "type": placeholder.placeholder_format.type.name if hasattr(placeholder.placeholder_format.type, 'name') else str(placeholder.placeholder_format.type),
                        "name": getattr(placeholder, "name", f"Placeholder {j}"),
                        "left": round(placeholder.left.inches, 2),
                        "top": round(placeholder.top.inches, 2),
                        "width": round(placeholder.width.inches, 2),
                        "height": round(placeholder.height.inches, 2)
                    }
                    layout_info["placeholders"].append(ph_info)
                except Exception as e:
                    layout_info["placeholders"].append({"error": str(e), "index": j})
            
            layouts_info.append(layout_info)
        
        return {
            "total_layouts": len(prs.slide_layouts),
            "layouts": layouts_info
        }
    except Exception as e:
        return {"error": str(e)}

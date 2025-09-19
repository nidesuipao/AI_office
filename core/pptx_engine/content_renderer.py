"""
内容渲染器 - 负责具体内容的渲染（文本、图片、表格）
"""

import os
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class ContentRenderer:
    """内容渲染器，负责文本、图片、表格的具体渲染"""
    
    def __init__(self, presentation: Presentation, font_calculator=None):
        self.prs = presentation
        self.font_calculator = font_calculator
        self.md_base_dir = None
    
    def set_md_base_dir(self, md_base_dir: str):
        """设置Markdown文件基础目录，用于解析图片相对路径"""
        self.md_base_dir = md_base_dir
    
    def _get_slide_index(self, slide) -> int:
        """根据 slide 对象获取其索引"""
        try:
            for idx, s in enumerate(self.prs.slides):
                if s == slide:
                    return idx
        except Exception:
            pass
        return -1
    
    def _resolve_image_path(self, src: str) -> str:
        """解析图片路径，支持相对路径和绝对路径"""
        try:
            if os.path.isabs(src):
                return src
            if self.md_base_dir:
                candidate = os.path.join(self.md_base_dir, src)
                if os.path.exists(candidate):
                    return candidate
            # 退一步：相对当前工作目录
            return src
        except Exception:
            return src
    
    def add_text_box(self, slide_index: int, text: str, left: float = 1.0, top: float = 1.0, 
                    width: float = 4.0, height: float = 1.0, font_name: Optional[str] = None, 
                    font_size: Optional[int] = None, font_bold: bool = False, 
                    font_italic: bool = False, text_color: Optional[str] = None, 
                    alignment: str = "left") -> Optional[object]:
        """添加文本框"""
        try:
            slide = self.prs.slides[slide_index]
            textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            frame = textbox.text_frame
            frame.clear()
            para = frame.paragraphs[0]
            run = para.add_run()
            run.text = text
            
            if font_name:
                run.font.name = font_name
            else:
                # 从FontCalculator获取默认字体名称
                run.font.name = '微软雅黑'  # 保持向后兼容
            if font_size:
                run.font.size = Pt(font_size)
            run.font.bold = bool(font_bold)
            run.font.italic = bool(font_italic)
            
            if text_color:
                hexv = text_color.lstrip('#')
                if len(hexv) == 6:
                    run.font.color.rgb = RGBColor(int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16))
            
            align_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
                'justify': PP_ALIGN.JUSTIFY
            }
            para.alignment = align_map.get(alignment, PP_ALIGN.LEFT)
            return textbox
        except Exception as e:
            print(f"添加文本框失败: {e}")
            return None
    
    def insert_image(self, slide_index: int, image_path: str, left: float = 1.0, top: float = 1.0, 
                    width: float = None, height: float = None, caption: str = None) -> Optional[object]:
        """插入图片，支持标题"""
        try:
            slide = self.prs.slides[slide_index]
            img_path = self._resolve_image_path(image_path)
            kwargs = {}
            if width is not None:
                kwargs['width'] = Inches(width)
            if height is not None:
                kwargs['height'] = Inches(height)
            pic = slide.shapes.add_picture(img_path, Inches(left), Inches(top), **kwargs)
            
            # 如果有标题，在图片下方添加标题
            if caption:
                # 计算标题位置：图片底部 + 0.1英寸间距
                caption_top = top + (height if height else 2.0) + 0.1
                # 确保标题不会超过页面底部
                page_bottom = 7.5  # 假设页面高度为7.5英寸
                if caption_top + 0.3 > page_bottom:  # 0.3英寸是标题高度
                    caption_top = page_bottom - 0.3
                
                # 添加标题文本框
                caption_textbox = slide.shapes.add_textbox(
                    Inches(left), Inches(caption_top), 
                    Inches(width if width else 4.0), Inches(0.3)
                )
                caption_frame = caption_textbox.text_frame
                caption_frame.clear()
                caption_para = caption_frame.paragraphs[0]
                caption_run = caption_para.add_run()
                caption_run.text = caption
                
                # 设置标题格式
                caption_run.font.name = '微软雅黑'
                # 使用字体计算器计算标题字体大小（使用合理的可用高度）
                caption_size = self.font_calculator.calculate_optimal_font_size(1.0, 1, 'caption')
                caption_run.font.size = Pt(caption_size)
                caption_run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)  # 灰色
                caption_para.alignment = PP_ALIGN.CENTER  # 居中对齐
            
            return pic
        except Exception as e:
            print(f"插入图片失败: {e}")
            return None
    
    def insert_table(self, slide_index: int, rows: int, cols: int, data: Optional[List[List[str]]] = None, 
                    left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 3.0, 
                    caption: str = None) -> Optional[object]:
        """插入表格，支持标题"""
        try:
            slide = self.prs.slides[slide_index]
            table_obj = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
            table = table_obj.table
            
            if data is not None:
                for r in range(rows):
                    for c in range(cols):
                        cell = table.cell(r, c)
                        text_val = ''
                        if data and r < len(data) and c < len(data[r]):
                            text_val = str(data[r][c])
                        cell.text = text_val
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = '微软雅黑'
                                run.font.size = Pt(14)
                            paragraph.alignment = PP_ALIGN.CENTER
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 如果有标题，在表格下方添加标题
            if caption:
                # 计算标题位置：表格底部 + 0.1英寸间距
                caption_top = top + height + 0.1
                # 确保标题不会超过页面底部
                page_bottom = 7.5  # 假设页面高度为7.5英寸
                if caption_top + 0.3 > page_bottom:  # 0.3英寸是标题高度
                    caption_top = page_bottom - 0.3
                
                # 添加标题文本框
                caption_textbox = slide.shapes.add_textbox(
                    Inches(left), Inches(caption_top), 
                    Inches(width), Inches(0.3)
                )
                caption_frame = caption_textbox.text_frame
                caption_frame.clear()
                caption_para = caption_frame.paragraphs[0]
                caption_run = caption_para.add_run()
                caption_run.text = caption
                
                # 设置标题格式
                caption_run.font.name = '微软雅黑'
                # 使用字体计算器计算标题字体大小（使用合理的可用高度）
                caption_size = self.font_calculator.calculate_optimal_font_size(1.0, 1, 'caption')
                caption_run.font.size = Pt(caption_size)
                caption_run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)  # 灰色
                caption_para.alignment = PP_ALIGN.CENTER  # 居中对齐
            
            return table_obj
        except Exception as e:
            print(f"插入表格失败: {e}")
            return None
    
    def add_text_content_left_aligned(self, slide, text_blocks, content_top, content_height, font_calculator):
        """添加左对齐的文字内容，自动换行并适应文本框宽度"""
        idx = self._get_slide_index(slide)
        text_textbox = self.add_text_box(idx, "", 1.0, content_top, 11.0, content_height)
        
        text_frame = text_textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True  # 启用自动换行
        
        def _wrap_text(text: str, max_chars_per_line: int = 42) -> str:
            # 简单基于字符数的折行，兼容中英文混排
            lines = []
            current = ''
            for ch in text:
                current += ch
                if len(current) >= max_chars_per_line and ch == ' ':
                    lines.append(current.rstrip())
                    current = ''
            if current:
                lines.append(current)
            return '\n'.join(lines)
        
        # 创建第一个段落
        first_para = text_frame.paragraphs[0]
        first_para_used = False
        
        for block in text_blocks:
            if block['type'] == 'list':
                for item in block['items']:
                    if not first_para_used:
                        para = first_para
                        first_para_used = True
                    else:
                        para = text_frame.add_paragraph()
                    run = para.add_run()
                    run.text = f"• {_wrap_text(item)}"
                    run.font.name = '微软雅黑'
                    font_size = font_calculator.calculate_optimal_font_size(content_height, len(text_blocks), 'text')
                    run.font.size = Pt(font_size)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
                    para.alignment = PP_ALIGN.LEFT
                    para.space_after = Pt(6)
            elif block['type'] == 'paragraph':
                if not first_para_used:
                    para = first_para
                    first_para_used = True
                else:
                    para = text_frame.add_paragraph()
                run = para.add_run()
                run.text = _wrap_text(block['text'])
                run.font.name = '微软雅黑'
                total_text_items = sum(len(block.get('items', [])) if block['type'] == 'list' 
                                     else 1 for block in text_blocks)
                font_size = font_calculator.calculate_optimal_font_size(content_height, total_text_items, 'text')
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
                para.alignment = PP_ALIGN.LEFT
                para.space_after = Pt(8)
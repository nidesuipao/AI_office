"""
幻灯片构建器 - 负责特定类型幻灯片的构建
"""

import re
import os
import sys
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 添加项目根目录到Python路径
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils.pptx_utils import (
    update_text_preserve_format,
    split_title_by_length,
    smart_update_toc_items
)
from .content_renderer import ContentRenderer
from .font_calculator import FontCalculator


class SlideBuilder:
    """幻灯片构建器，负责特定类型幻灯片的构建"""
    
    def __init__(self, presentation: Presentation, content_renderer: ContentRenderer, font_calculator: FontCalculator):
        self.prs = presentation
        self.renderer = content_renderer
        self.font_calc = font_calculator
        
        # 存储章节信息，用于动态获取章节标题
        self.chapters_info = {}  # {chapter_number: chapter_title}
    
    def add_title_slide(self, data: dict) -> int:
        """更新标题页"""
        title = data["title_page"].get("title", "")
        org = data["title_page"].get("org", "")
        date = data["title_page"].get("date", "")
        # 仅修改模板第一页，不新增
        slide = self.prs.slides[0]

        # 根据模板结构直接定位：
        # 形状0: 标题 ("第一行标题\n第二行标题")
        # 形状1: 单位 ("xxxx有限公司") 
        # 形状2: 日期 ("xxxx年xx月xx日")
        
        try:
            if len(slide.shapes) >= 3:
                # 标题文本框 (形状0) - 自动拆分为两行
                if title:
                    formatted_title = split_title_by_length(title)
                    print(formatted_title)
                    update_text_preserve_format(slide.shapes[0], formatted_title)
                
                # 单位文本框 (形状1)  
                if org:
                    update_text_preserve_format(slide.shapes[1], org)
                
                # 日期文本框 (形状2)
                if date:
                    update_text_preserve_format(slide.shapes[2], date)

        except Exception as e:
            print(f"修改标题页文本时出错: {e}")
        
        return 1
    
    def add_toc_slide(self, data: dict) -> int:
        """更新模板第二页的目录内容，智能匹配2-5个目录项"""
        if len(self.prs.slides) < 2:
            raise ValueError("模板中不存在第二页，无法修改目录页")
        
        slide = self.prs.slides[1]  # 第二页
        toc_items = data.get("toc", [])
        
        try:
            # 使用智能目录更新功能，保持模板美观性
            smart_update_toc_items(slide, toc_items)
                    
        except Exception as e:
            print(f"修改目录页时出错: {e}")
        
        return 2
    
    def add_subsection_slides(self, subsections: list) -> int:
        """为每个子章节(### 1.1, 1.2等)添加单独的页面"""
        if not subsections:
            return 0
        
        # 使用布局1："标题和内容"
        layout = self.prs.slide_layouts[1]
        slides_added = 0
        
        current_chapter = 0
        
        for subsection in subsections:
            try:
                # 检查是否需要添加章节分隔页
                if subsection['chapter_number'] != current_chapter:
                    current_chapter = subsection['chapter_number']
                    # 从子章节信息中提取章节标题，并存储到chapters_info中
                    if 'chapter_title' in subsection:
                        self.chapters_info[current_chapter] = subsection['chapter_title']
                    chapter_title = self._get_chapter_title(current_chapter)
                    # 添加章节分隔页
                    divider_added = self.add_chapter_divider_slide(current_chapter, chapter_title)
                    if divider_added:
                        slides_added += 1
                
                # 添加子章节内容页
                slide = self.prs.slides.add_slide(layout)
                slides_added += 1
                
                # 去掉子章节标题中的编号（如"1.1 办公痛点" -> "办公痛点驱动自动化需求"）
                clean_subsection_title = re.sub(r'^\d+\.\d+\s*', '', subsection['title'])
                
                # 在蓝色背景区域添加子章节标题
                subsection_title_textbox = slide.shapes.add_textbox(
                    Inches(1.0),    # 左边距
                    Inches(0.15),   # 顶部位置，在蓝色背景条内
                    Inches(11.0),   # 宽度，几乎占满
                    Inches(0.5)     # 高度
                )
                
                # 设置子章节标题内容和格式
                title_frame = subsection_title_textbox.text_frame
                title_frame.clear()
                title_para = title_frame.paragraphs[0]
                title_run = title_para.add_run()
                title_run.text = clean_subsection_title
                
                # 设置标题格式：白色字体，居中，加粗
                title_run.font.name = '微软雅黑'
                # 根据标题区域高度自动调整字体大小
                title_font_size = self.font_calc.calculate_title_font_size(0.5)  # 标题区域固定0.5英寸
                title_run.font.size = Pt(title_font_size)
                title_run.font.bold = True
                title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # 白色字体
                title_para.alignment = PP_ALIGN.CENTER  # 居中对齐
                
                print(f"已添加子章节标题: {clean_subsection_title} (原标题: {subsection['title']})")
                
                # 创建内容区域（在子章节标题下方）
                content_top = 1.0   # 在标题下方开始
                content_height = 5.8  # 内容区域高度
                content_textbox = slide.shapes.add_textbox(
                    Inches(1.0),    # 左边距
                    Inches(content_top),     # 内容顶部
                    Inches(11.0),   # 宽度
                    Inches(content_height)   # 高度
                )
                
                # 智能布局自动匹配
                from .layout_manager import LayoutManager
                layout_manager = LayoutManager(self.renderer, self.font_calc)
                layout_manager.add_content_auto_layout(slide, subsection['content_blocks'], content_top, content_height)
                
                print(f"已添加子章节: {clean_subsection_title}")
                
            except Exception as e:
                print(f"添加子章节'{subsection['title']}'时出错: {e}")
        
        return slides_added
    
    def add_chapter_divider_slide(self, chapter_number: int, chapter_title: str) -> int:
        """添加章节分隔页，使用布局2 '2_标题幻灯片'，包含'0X'和'章节标题'占位符"""
        try:
            # 使用布局2："2_标题幻灯片"，现在包含"0X"和"目录标题"占位符
            layout2 = self.prs.slide_layouts[2]
            
            # 添加章节分隔幻灯片
            divider_slide = self.prs.slides.add_slide(layout2)
            
            # 更新占位符内容（根据占位符类型和索引）
            ox_updated = False
            title_updated = False
            
            for shape in divider_slide.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        ph_type = shape.placeholder_format.type.name
                        ph_idx = shape.placeholder_format.idx
                        
                        # 根据模板分析：
                        # - "0X"是 BODY 类型，索引14 的占位符
                        # - "目录标题"是 OBJECT 类型，索引13 的占位符
                        if ph_type == 'BODY' and ph_idx == 14:
                            # 更新章节编号并确保居中对齐
                            shape.text = f"{chapter_number:02d}"
                            # 强制设置居中对齐
                            if shape.text_frame and shape.text_frame.paragraphs:
                                para = shape.text_frame.paragraphs[0]
                                para.alignment = PP_ALIGN.CENTER
                            ox_updated = True
                            print(f"  更新章节编号占位符 (BODY-14): {chapter_number:02d} [居中对齐]")
                        
                        elif ph_type == 'OBJECT' and ph_idx == 13:
                            # 去掉章节标题中的编号（如"1. "、"2. "等）
                            clean_title = re.sub(r'^\d+\.\s*', '', chapter_title)
                            # 更新章节标题并保持原有格式
                            update_text_preserve_format(shape, clean_title)
                            title_updated = True
                            print(f"  更新章节标题占位符 (OBJECT-13): {clean_title} (原标题: {chapter_title})")
                        
                    except Exception as e:
                        print(f"  处理占位符时出错: {e}")
            
            # 删除其他不需要的占位符（保留0X和标题占位符）
            shapes_to_remove = []
            for shape in divider_slide.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        # 检查占位符类型，只保留我们需要的
                        if hasattr(shape, 'placeholder_format'):
                            ph_type = shape.placeholder_format.type.name
                            # 删除标题、副标题、日期、页脚、页码等占位符，保留内容相关的
                            if ph_type in ['CENTER_TITLE', 'SUBTITLE', 'DATE', 'FOOTER', 'SLIDE_NUMBER']:
                                shapes_to_remove.append(shape)
                    except:
                        # 如果无法确定占位符类型，检查文本内容
                        if hasattr(shape, 'text') and shape.text:
                            if '0X' not in shape.text and '标题' not in shape.text:
                                if any(word in shape.text for word in ['单击此处', '编辑']):
                                    shapes_to_remove.append(shape)
            
            # 删除不需要的占位符
            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)
            
            print(f"已添加章节{chapter_number:02d}分隔页: {chapter_title} (占位符更新: 编号={ox_updated}, 标题={title_updated})")
            return 1
                
        except Exception as e:
            print(f"添加章节分隔页时出错: {e}")
            return 0
    
    def add_thanks_slide(self) -> int:
        """添加致谢页，直接导入布局3 '1_标题幻灯片'模板的所有元素"""
        try:
            # 使用布局3："1_标题幻灯片"
            layout3 = self.prs.slide_layouts[3]
            
            # 添加致谢幻灯片
            thanks_slide = self.prs.slides.add_slide(layout3)
            
            # 删除所有占位符，只保留背景和"感 谢 聆 听"文本
            shapes_to_remove = []
            for shape in thanks_slide.shapes:
                # 找到所有占位符并标记删除
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    shapes_to_remove.append(shape)
            
            # 删除所有占位符
            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)
            
            print(f"已删除 {len(shapes_to_remove)} 个占位符")
            
            print("已添加致谢页（使用标题占位符）")
            return 1
                
        except Exception as e:
            print(f"添加致谢页时出错: {e}")
            return 0
    
    def _get_chapter_title(self, chapter_number: int) -> str:
        """根据章节编号获取章节标题，优先从已解析的章节信息中获取"""
        # 首先从self.chapters_info中获取（如果从md文件解析过）
        if chapter_number in self.chapters_info:
            return self.chapters_info[chapter_number]
        
        # 如果没有解析过，使用默认格式
        return f"{chapter_number}. 章节标题"
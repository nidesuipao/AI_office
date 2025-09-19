"""
PPTX构建器主类 - 负责整体协调和公共接口
"""

import os
import sys
from typing import List, Optional

# 添加项目根目录到Python路径
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# 导入工具函数
from utils.pptx_utils import (
    parse_md_for_ppt_structure,
    read_text,
    extract_subsection_content_from_md
)

# 导入组件
from .font_calculator import FontCalculator
from .content_renderer import ContentRenderer
from .layout_manager import LayoutManager
from .slide_builder import SlideBuilder

pptx_installed = True
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.oxml import parse_xml
except ImportError:
    print("警告: 未检测到python-pptx库，功能将不可用 (pip install python-pptx)")
    pptx_installed = False


class PPTXBuilder:
    """PPTX构建器主类，负责整体协调和公共接口"""
    
    # 布局索引常量
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_CHAPTER_DIVIDER = 2
    
    # 尺寸常量（英寸）
    SUBSECTION_TITLE_LEFT = 1.0
    SUBSECTION_TITLE_TOP = 0.15
    SUBSECTION_TITLE_WIDTH = 11.0
    SUBSECTION_TITLE_HEIGHT = 0.5
    CONTENT_LEFT = 1.0
    CONTENT_HEIGHT = 5.8
    
    # 字体格式常量
    FONT_NAME = '微软雅黑'
    TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)  # 白色
    
    def __init__(self, file_path: Optional[str] = None, template_path: Optional[str] = None):
        if not pptx_installed:
            raise RuntimeError("python-pptx 库未安装，请先安装: pip install python-pptx")
        
        # 基础属性
        self.file_path = file_path
        self.prs = None
        self.md_base_dir = None
        
        # 初始化组件
        self._init_components(file_path, template_path)
    
    def _init_components(self, file_path: Optional[str], template_path: Optional[str]):
        """初始化各个组件"""
        # 初始化Presentation
        if file_path and os.path.exists(file_path):
            self.prs = Presentation(file_path)
        elif template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # 初始化各个组件
        self.font_calc = FontCalculator()
        self.renderer = ContentRenderer(self.prs, self.font_calc)
        self.layout_manager = LayoutManager(self.renderer, self.font_calc)
        self.slide_builder = SlideBuilder(self.prs, self.renderer, self.font_calc)
    
    def _reinit_components(self):
        """重新初始化组件，使用当前的Presentation对象"""
        self.renderer = ContentRenderer(self.prs, self.font_calc)
        self.layout_manager = LayoutManager(self.renderer, self.font_calc)
        self.slide_builder = SlideBuilder(self.prs, self.renderer, self.font_calc)
    
    # 基础操作
    def save(self, output_path: Optional[str] = None) -> str:
        """保存PPT文件"""
        path = output_path or self.file_path
        if not path:
            raise ValueError("未指定保存路径")
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)
        return path

    def info(self) -> str:
        """获取PPT信息"""
        slide_count = len(self.prs.slides)
        lines = [f"幻灯片数量: {slide_count}", "", "幻灯片概览:"]
        for i, slide in enumerate(self.prs.slides):
            title = "无标题"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and (getattr(shape, "name", "").startswith("Title") or "标题" in getattr(shape, "name", "")):
                    txt = shape.text
                    title = txt[:50] + "..." if len(txt) > 50 else txt
                    break
            lines.append(f"幻灯片 {i+1}: {title} (形状 {len(slide.shapes)})")
        return "\n".join(lines)
    
    # 幻灯片操作
    def add_slide(self, layout_name: str = "Title and Content") -> int:
        """添加幻灯片"""
        try:
            chosen_layout = None
            for layout in self.prs.slide_layouts:
                try:
                    if getattr(layout, 'name', '') == layout_name:
                        chosen_layout = layout
                        break
                except Exception:
                    continue
            if chosen_layout is None:
                # 常用默认布局：标题和内容，一般索引为1
                fallback_idx = 1 if len(self.prs.slide_layouts) > 1 else 0
                chosen_layout = self.prs.slide_layouts[fallback_idx]
            self.prs.slides.add_slide(chosen_layout)
            # 返回新页的1-based索引
            return len(self.prs.slides)
        except Exception as e:
            print(f"添加幻灯片失败: {e}")
            return -1
    
    # 内容操作（委托给ContentRenderer）
    def add_text_box(self, slide_index: int, text: str, left: float = 1.0, top: float = 1.0, 
                    width: float = 4.0, height: float = 1.0, font_name: Optional[str] = None, 
                    font_size: Optional[int] = None, font_bold: bool = False, 
                    font_italic: bool = False, text_color: Optional[str] = None, 
                    alignment: str = "left") -> Optional[object]:
        """添加文本框"""
        return self.renderer.add_text_box(slide_index, text, left, top, width, height, 
                                        font_name, font_size, font_bold, font_italic, 
                                        text_color, alignment)
    
    def insert_image(self, slide_index: int, image_path: str, left: float = 1.0, top: float = 1.0, 
                    width: float = None, height: float = None, caption: str = None) -> Optional[object]:
        """插入图片"""
        return self.renderer.insert_image(slide_index, image_path, left, top, width, height, caption)
    
    def insert_table(self, slide_index: int, rows: int, cols: int, data: Optional[List[List[str]]] = None, 
                    left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 3.0, 
                    caption: str = None) -> Optional[object]:
        """插入表格"""
        return self.renderer.insert_table(slide_index, rows, cols, data, left, top, width, height, caption)
    
    # 幻灯片构建操作（委托给SlideBuilder）
    def add_title_slide(self, data: dict) -> int:
        """添加标题页"""
        return self.slide_builder.add_title_slide(data)
    
    def add_toc_slide(self, data: dict) -> int:
        """添加目录页"""
        return self.slide_builder.add_toc_slide(data)
    
    def add_subsection_slides(self, subsections: list) -> int:
        """添加子章节页面"""
        return self.slide_builder.add_subsection_slides(subsections)
    
    def add_chapter_divider_slide(self, chapter_number: int, chapter_title: str) -> int:
        """添加章节分隔页"""
        return self.slide_builder.add_chapter_divider_slide(chapter_number, chapter_title)
    
    def add_thanks_slide(self) -> int:
        """添加致谢页"""
        return self.slide_builder.add_thanks_slide()
    
    # 字体计算操作（委托给FontCalculator）
    def calculate_optimal_font_size(self, available_height, content_amount, content_type):
        """计算最佳字体大小"""
        return self.font_calc.calculate_optimal_font_size(available_height, content_amount, content_type)
    
    def calculate_title_font_size(self, available_height):
        """计算标题字体大小"""
        return self.font_calc.calculate_title_font_size(available_height)
    
    def calculate_table_font_size(self, table_height, rows, cols, content_type):
        """计算表格字体大小"""
        return self.font_calc.calculate_table_font_size(table_height, rows, cols, content_type)
    
    # from-md 主要功能
    def from_md(self, md_path: str, template_path: str, output_path: str) -> str:
        """从Markdown文件生成PPT"""
        # 解析 Markdown，并在不增加页数的前提下修改模板第一页，保存到输出
        if not os.path.exists(md_path):
            raise FileNotFoundError(f"Markdown 文件不存在: {md_path}")
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"模板不存在: {template_path}")
        
        md_text = read_text(md_path)
        
        # 记录 Markdown 所在目录，供图片相对路径解析
        try:
            self.md_base_dir = os.path.dirname(md_path)
            self.renderer.set_md_base_dir(self.md_base_dir)
        except Exception:
            self.md_base_dir = None
        
        data = parse_md_for_ppt_structure(md_text)
        print(data)

        # 使用模板初始化，并修改第一页和目录页
        self.prs = Presentation(template_path)
        # 重新初始化组件以使用新的presentation
        self._reinit_components()
        if self.md_base_dir:
            self.renderer.set_md_base_dir(self.md_base_dir)
        
        if len(self.prs.slides) == 0:
            raise ValueError("模板中不存在任何幻灯片，无法修改第一页")
        
        # 更新标题页
        self.add_title_slide(data)
        
        # 更新目录页（如果存在第二页）
        if len(self.prs.slides) >= 2:
            self.add_toc_slide(data)
        else:
            print("警告: 模板中不存在第二页，跳过目录页更新")

        # 添加子章节内容页面（每个1.1、1.2等子标题单独成页）
        subsections = extract_subsection_content_from_md(md_text)
        if subsections:
            slides_added = self.add_subsection_slides(subsections)
            print(f"已添加 {slides_added} 个子章节幻灯片（包含章节分隔页）")
        else:
            print("警告: 未找到子章节内容，跳过子章节页面生成")

        # 添加致谢页
        thanks_added = self.add_thanks_slide()
        if thanks_added:
            print("已在最后添加致谢页")

        # 保存到输出
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        return output_path
